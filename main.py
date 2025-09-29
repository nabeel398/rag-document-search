from fastapi import FastAPI, UploadFile, File, Form
from typing import List
import os
import shutil
import tempfile

from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
from langchain.chains import RetrievalQA
from langchain_groq import ChatGroq  # Groq integration
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.embeddings import HuggingFaceEmbeddings

from PyPDF2 import PdfReader
import docx
import pptx
from fastapi.middleware.cors import CORSMiddleware

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # You can restrict this to ["http://localhost:5500"] etc.
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Sentence Transformer embedding model
embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
FAISS_INDEX_PATH = "faiss_index"

# Initialize Groq LLM
llm = ChatGroq(
    groq_api_key=os.getenv("GROQ_API_KEY"),
    model_name="llama-3.3-70b-versatile",  # Updated model
    temperature=0.2
)

# Load or initialize FAISS
vector_store = None
if os.path.exists(FAISS_INDEX_PATH):
    vector_store = FAISS.load_local(FAISS_INDEX_PATH, embedding_model, allow_dangerous_deserialization=True)

# File reader function
def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        reader = PdfReader(file_path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    elif ext == ".docx":
        doc_file = docx.Document(file_path)
        return "\n".join([p.text for p in doc_file.paragraphs])
    elif ext == ".pptx":
        ppt = pptx.Presentation(file_path)
        text = [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")]
        return "\n".join(text)
    elif ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    return ""

# Upload documents and embed them
@app.post("/upload")
async def upload_files(files: List[UploadFile] = File(...)):
    global vector_store
    docs = []
    with tempfile.TemporaryDirectory() as tmpdir:
        for file in files:
            temp_path = os.path.join(tmpdir, file.filename)
            with open(temp_path, "wb") as f:
                shutil.copyfileobj(file.file, f)
            content = extract_text(temp_path)
            if content.strip():
                docs.append(Document(page_content=content))

    # Chunking and Overlap
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    chunks = []
    for doc in docs:
        chunks.extend([Document(page_content=chunk) for chunk in text_splitter.split_text(doc.page_content)])

    if chunks:
        if vector_store is None:
            vector_store = FAISS.from_documents(chunks, embedding_model)
        else:
            vector_store.add_documents(chunks)
        vector_store.save_local(FAISS_INDEX_PATH)
    return {"status": "success", "documents_added": len(chunks)}

# Query from stored embeddings
@app.post("/query")
async def query_vectorstore(query: str = Form(...)):
    if vector_store is None:
        return {"error": "No documents uploaded yet."}

    retriever = vector_store.as_retriever(search_kwargs={"k": 3})

    qa_chain = RetrievalQA.from_chain_type(
        llm=llm,
        retriever=retriever,
        chain_type="stuff"
    )

    result = qa_chain.run(query)

    # ✅ Check if answer is empty or irrelevant
    if not result or result.strip() == "" or "I don't know" in result:
        return {
            "query": query,
            "answer": "answer not found, try some other question"
        }

    return {
        "query": query,
        "answer": result
    }
